# tests/test_ppt_report.py
"""
Test scaffolds for REP-02 requirement.
Initially skipped - will pass after src/report/ppt_report.py implementation.
"""

import pytest
from pathlib import Path
import tempfile
from pptx import Presentation
from pptx.presentation import Presentation as PresentationClass


class TestPPTReport:
    """Tests for PPT report generation via python-pptx."""

    def test_generate_ppt(self, tmp_path):
        """REP-02: Verify PPTReportGenerator creates Presentation object with style."""
        from src.report.ppt_report import PPTReportGenerator
        from pptx import Presentation

        # Test initialization with style_id
        generator = PPTReportGenerator(style_id='ft')

        # Verify Presentation object created
        assert generator.prs is not None
        assert isinstance(generator.prs, PresentationClass)

        # Verify style loaded
        assert generator.style is not None
        assert generator.style.id == 'ft'

        # Test basic save
        output_path = tmp_path / 'test.pptx'
        generator.save_ppt(output_path)

        # Verify PPT file created
        assert output_path.exists()
        assert output_path.suffix == '.pptx'

        # Verify file size is reasonable (not empty)
        assert output_path.stat().st_size > 1000  # PPTX should be at least 1KB

    def test_ppt_slide_structure(self, tmp_path):
        """Verify PPT slides have title and content structure."""
        from src.report.ppt_report import PPTReportGenerator
        from pptx import Presentation

        generator = PPTReportGenerator(style_id='ft')

        # Add slide with title and content
        generator.add_slide('测试标题', '测试内容')

        # Save to verify structure
        output_path = tmp_path / 'test_slide.pptx'
        generator.save_ppt(output_path)

        # Open and verify structure
        prs = Presentation(str(output_path))

        # Verify slide exists
        assert len(prs.slides) >= 1

        # Verify first slide has title
        first_slide = prs.slides[0]
        assert len(first_slide.shapes) > 0

        # Verify title text
        title_shape = first_slide.shapes.title
        assert title_shape is not None
        assert title_shape.text == '测试标题'

        # Verify content placeholder has text
        placeholder = first_slide.placeholders[1]
        assert placeholder.text == '测试内容'

    def test_ppt_chart_embedding(self, tmp_path):
        """Verify charts embedded as images in PPT."""
        from src.report.ppt_report import PPTReportGenerator
        from pptx import Presentation
        from pptx.shapes.picture import Picture
        import plotly.graph_objects as go
        import pandas as pd

        generator = PPTReportGenerator(style_id='ft')

        # Create a simple Plotly chart
        df = pd.DataFrame({'x': ['A', 'B', 'C'], 'y': [10, 20, 15]})
        fig = go.Figure()
        fig.add_trace(go.Bar(x=df['x'], y=df['y']))

        # Add chart slide
        generator.add_chart_slide('图表测试', fig, tmp_path)

        # Save to verify
        output_path = tmp_path / 'test_chart.pptx'
        generator.save_ppt(output_path)

        # Open and verify chart images
        prs = Presentation(str(output_path))

        # Verify slide exists
        assert len(prs.slides) >= 1

        # Check for picture shapes (chart images)
        picture_count = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'image') or isinstance(shape, Picture):
                    picture_count += 1

        # Should have at least one chart image
        assert picture_count >= 1

    def test_build_report_workflow(self, tmp_path):
        """Verify build_report creates multi-slide presentation."""
        from src.report.ppt_report import PPTReportGenerator
        from pptx import Presentation
        import plotly.graph_objects as go
        import pandas as pd

        generator = PPTReportGenerator(style_id='mckinsey')

        # Create synthesis data
        synthesis_data = {
            'report_title': '数据分析报告',
            'executive_summary': '本次分析发现ROI增长15%，关键指标表现良好。',
            'metrics_summary': 'ROI: 15.5%, CTR: 3.2%, 转化率: 5.8%',
            'conclusions_text': '建议继续优化投放策略，重点关注转化率提升。'
        }

        # Create sample charts
        df = pd.DataFrame({'x': ['A', 'B', 'C'], 'y': [10, 20, 15]})
        fig1 = go.Figure()
        fig1.add_trace(go.Bar(x=df['x'], y=df['y']))

        fig2 = go.Figure()
        fig2.add_trace(go.Scatter(x=df['x'], y=df['y'], mode='lines+markers'))

        charts = [fig1, fig2]

        # Build report
        generator.build_report(synthesis_data, charts, tmp_path)

        # Save
        output_path = tmp_path / 'full_report.pptx'
        generator.save_ppt(output_path)

        # Verify file created
        assert output_path.exists()
        assert output_path.stat().st_size > 5000  # Multi-slide PPTX should be larger

        # Open and verify slide count
        prs = Presentation(str(output_path))

        # Should have: title + metrics + 2 charts + conclusions = 5 slides minimum
        assert len(prs.slides) >= 5