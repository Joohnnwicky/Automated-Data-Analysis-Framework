"""PPT report generation using python-pptx direct construction.

This module implements REP-02: PPT generation for presentation scenarios
when user explicitly requests PPT output.

Uses python-pptx to build slides directly from synthesis data (no HTML conversion).
Charts are embedded as PNG images exported from Plotly.

Threat Model:
- T-4-02: python-pptx handles text safely (no XSS in PPT)
- T-4-03: PNG export safe (binary format, no data injection)
"""

from pathlib import Path
from typing import List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import plotly.graph_objects as go
import plotly.io as pio

from src.report.styles import get_design_style, DesignStyle


def hex_to_rgb(hex_str: str) -> RGBColor:
    """Convert hex color string to RGBColor.

    Args:
        hex_str: Hex color string (e.g., '#E3120B' or 'E3120B')

    Returns:
        RGBColor object for python-pptx.

    Example:
        >>> hex_to_rgb('#E3120B')
        RGBColor(227, 18, 11)
    """
    # Strip '#' prefix if present
    hex_str = hex_str.lstrip('#')

    # Parse hex to RGB values
    r = int(hex_str[0:2], 16)
    g = int(hex_str[2:4], 16)
    b = int(hex_str[4:6], 16)

    return RGBColor(r, g, b)


class PPTReportGenerator:
    """Generate PPT report using python-pptx direct construction.

    REP-02: Creates PowerPoint slides directly from synthesis data.
    No HTML→PPTX conversion - builds slides programmatically.

    Attributes:
        style: DesignStyle object for visual formatting.
        prs: python-pptx Presentation object.
        primary_color_rgb: RGBColor for title fonts.

    Example:
        >>> generator = PPTReportGenerator(style_id='ft')
        >>> generator.add_slide('Executive Summary', 'Key findings...')
        >>> generator.save_ppt(Path('output/report.pptx'))
    """

    def __init__(self, style_id: str = 'ft'):
        """Initialize PPTReportGenerator with design style.

        Args:
            style_id: Design style identifier (e.g., 'ft', 'mckinsey').
                      Defaults to 'ft' (Financial Times style).

        Creates:
            - Presentation object (blank)
            - DesignStyle loaded via get_design_style()
            - RGBColor for primary color
        """
        self.style = get_design_style(style_id)
        self.prs = Presentation()
        self.primary_color_rgb = hex_to_rgb(self.style.primary_color)

    def add_slide(self, title: str, content: Optional[str] = None):
        """Add slide with title and optional content.

        Uses slide_layouts[1] (Title and Content layout).

        Args:
            title: Slide title text.
            content: Optional content text for placeholder.

        Note:
            Title font size is Pt(24), color is primary_color_rgb.
        """
        # Use Title and Content layout (index 1)
        layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(layout)

        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title

        # Set title font style
        for paragraph in title_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(24)
                run.font.color.rgb = self.primary_color_rgb

        # Set content if provided
        if content:
            # Placeholder index 1 is the content area
            placeholder = slide.placeholders[1]
            placeholder.text = content

    def add_chart_slide(self, title: str, chart_fig: go.Figure, tmp_path: Path):
        """Add slide with embedded chart image.

        Uses slide_layouts[5] (Blank layout) for flexible positioning.

        Args:
            title: Slide title text.
            chart_fig: Plotly Figure object to embed.
            tmp_path: Path for temporary PNG file storage.

        Note:
            Chart is exported as PNG (900x500) via kaleido.
            Chart positioned at Inches(0.5, 1.5) with 9x5 inch dimensions.
        """
        # Use Blank layout (index 5)
        layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(layout)

        # Add title textbox
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(9), Inches(0.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = title

        # Set title font color
        for paragraph in title_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(20)
                run.font.color.rgb = self.primary_color_rgb

        # Export chart as PNG
        png_bytes = pio.to_image(chart_fig, format='png', width=900, height=500)

        # Save to temporary file
        chart_png = tmp_path / 'chart.png'
        chart_png.write_bytes(png_bytes)

        # Add chart image to slide
        slide.shapes.add_picture(
            str(chart_png),
            Inches(0.5), Inches(1.5),
            width=Inches(9), height=Inches(5)
        )

    def save_ppt(self, output_path: Path) -> Path:
        """Save PPT file to output path.

        Args:
            output_path: Path to save .pptx file.

        Returns:
            Path to saved file.
        """
        self.prs.save(str(output_path))
        return output_path

    def build_report(
        self,
        synthesis_data: dict,
        charts: List[go.Figure],
        tmp_path: Path
    ):
        """Build complete PPT report from synthesis data and charts.

        Creates slides:
        1. Title slide with report title and executive summary
        2. Metrics slide with key indicators
        3. Chart slides (one per chart)
        4. Conclusions slide

        Args:
            synthesis_data: Dict with report_title, executive_summary,
                           metrics_summary, conclusions_text.
            charts: List of Plotly Figure objects.
            tmp_path: Path for temporary PNG file storage.
        """
        # Title slide
        self.add_slide(
            synthesis_data.get('report_title', '数据分析报告'),
            synthesis_data.get('executive_summary', '')
        )

        # Metrics slide
        self.add_slide(
            '关键指标',
            synthesis_data.get('metrics_summary', '')
        )

        # Chart slides
        chart_titles = [
            '趋势分析', '分类对比', '分布分析', '相关性分析'
        ]
        for i, chart in enumerate(charts):
            title = chart_titles[i] if i < len(chart_titles) else f'图表 {i+1}'
            self.add_chart_slide(title, chart, tmp_path)

        # Conclusions slide
        self.add_slide(
            '综合结论',
            synthesis_data.get('conclusions_text', '')
        )


__all__ = ['PPTReportGenerator', 'hex_to_rgb']